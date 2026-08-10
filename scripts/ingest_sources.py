#!/usr/bin/env python3
"""Ingest VaaniQ source PDFs/PPTX into citable markdown.

Phase 0 Step 0: convert Capstone proposal and Topic Approval deck into
machine-checkable markdown with page/slide markers and an ingest report.

Usage:
    python scripts/ingest_sources.py
    python scripts/ingest_sources.py --proposal PATH --approval PATH
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import fitz  # pymupdf

ROOT = Path(__file__).resolve().parents[1]
DEFAULT_PROPOSAL = Path(
    r"c:\Users\Aarav Phutane\Downloads\Capstone_Project_Proposal_260806_145819.pdf"
)
DEFAULT_APPROVAL = Path(
    r"c:\Users\Aarav Phutane\Downloads"
    r"\VaaniQ_Topic_Approval_Presentation_Final_B091_B093_B094_B106 (1).pdf"
)
OUT_DIR = ROOT / "docs" / "source"
FIGURES_DIR = OUT_DIR / "figures"
LOW_CHAR_THRESHOLD = 50


def _clean_text(text: str) -> str:
    """Normalize extracted PDF text for readable markdown."""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # Collapse runs of 3+ blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _page_has_images(page: fitz.Page) -> bool:
    return bool(page.get_images(full=True))


def extract_pdf(
    pdf_path: Path,
    *,
    marker_kind: str,
    render_low_char: bool,
    figure_prefix: str,
) -> tuple[str, list[dict[str, Any]]]:
    """Extract a PDF into markdown with page/slide markers.

    Args:
        pdf_path: Path to the PDF.
        marker_kind: ``page`` or ``slide`` for HTML comment markers.
        render_low_char: If True, rasterize pages with <50 chars of text.
        figure_prefix: Filename prefix for rendered PNGs.

    Returns:
        Tuple of (markdown body, per-page report dicts).
    """
    doc = fitz.open(pdf_path)
    parts: list[str] = [
        f"# Source extract: `{pdf_path.name}`",
        "",
        f"- Ingested: `{datetime.now(timezone.utc).isoformat()}`",
        f"- Source path: `{pdf_path}`",
        f"- Page/slide count: {doc.page_count}",
        "",
        "---",
        "",
    ]
    page_reports: list[dict[str, Any]] = []

    for i, page in enumerate(doc):
        n = i + 1
        raw = page.get_text("text")
        text = _clean_text(raw)
        char_count = len(text)
        image_count = len(page.get_images(full=True))
        low_char = char_count < LOW_CHAR_THRESHOLD
        rendered_path: str | None = None

        parts.append(f"<!-- {marker_kind}: {n} -->")
        parts.append("")
        parts.append(f"## {marker_kind.capitalize()} {n}")
        parts.append("")

        if text:
            parts.append(text)
            parts.append("")
        else:
            parts.append("*(no extractable text)*")
            parts.append("")

        if image_count:
            for img_i in range(image_count):
                parts.append(
                    f"<!-- figure: {marker_kind} {n}, embedded image {img_i + 1} "
                    f"of {image_count} -->"
                )
            parts.append("")

        if low_char and render_low_char:
            FIGURES_DIR.mkdir(parents=True, exist_ok=True)
            out_png = FIGURES_DIR / f"{figure_prefix}_{n:03d}.png"
            # 2x zoom for readability
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            pix.save(out_png)
            rendered_path = str(out_png.relative_to(ROOT)).replace("\\", "/")
            parts.append(
                f"![Low-text {marker_kind} {n}]({Path('figures') / out_png.name})"
            )
            parts.append("")
            parts.append(
                f"> **Manual review required:** {marker_kind} {n} yielded "
                f"{char_count} chars (< {LOW_CHAR_THRESHOLD}). "
                f"Rasterized to `{rendered_path}`."
            )
            parts.append("")

        # Attempt table-ish block extraction via text dict (best-effort)
        # pymupdf does not always recover true tables; leave raw text as source.

        page_reports.append(
            {
                f"{marker_kind}_number": n,
                "char_count": char_count,
                "image_count": image_count,
                "low_char": low_char,
                "rendered_figure": rendered_path,
            }
        )

    doc.close()
    return "\n".join(parts).rstrip() + "\n", page_reports


def try_pptx_notes(pptx_path: Path) -> str | None:
    """If a PPTX twin exists, extract speaker notes per slide."""
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError:
        return None

    if not pptx_path.is_file():
        return None

    prs = Presentation(str(pptx_path))
    lines: list[str] = [
        "",
        "---",
        "",
        "# Speaker notes (from PPTX twin)",
        "",
    ]
    for i, slide in enumerate(prs.slides, start=1):
        notes_frame = slide.notes_slide.notes_text_frame if slide.has_notes_slide else None
        notes = (notes_frame.text or "").strip() if notes_frame else ""
        lines.append(f"<!-- slide: {i} notes -->")
        lines.append("")
        lines.append(f"### Notes (slide {i})")
        lines.append("")
        lines.append(notes if notes else "*(no speaker notes)*")
        lines.append("")
    return "\n".join(lines)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--proposal", type=Path, default=DEFAULT_PROPOSAL)
    parser.add_argument("--approval", type=Path, default=DEFAULT_APPROVAL)
    parser.add_argument(
        "--approval-pptx",
        type=Path,
        default=None,
        help="Optional PPTX twin for speaker notes",
    )
    args = parser.parse_args()

    if not args.proposal.is_file():
        print(f"ERROR: proposal not found: {args.proposal}", file=sys.stderr)
        return 1
    if not args.approval.is_file():
        print(f"ERROR: approval deck not found: {args.approval}", file=sys.stderr)
        return 1

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    FIGURES_DIR.mkdir(parents=True, exist_ok=True)

    proposal_md, proposal_pages = extract_pdf(
        args.proposal,
        marker_kind="page",
        render_low_char=True,
        figure_prefix="proposal_page",
    )
    proposal_out = OUT_DIR / "Capstone_Project_Proposal.md"
    proposal_out.write_text(proposal_md, encoding="utf-8")

    approval_md, approval_slides = extract_pdf(
        args.approval,
        marker_kind="slide",
        render_low_char=True,
        figure_prefix="approval_slide",
    )

    pptx_twin = args.approval_pptx
    if pptx_twin is None:
        # Heuristic: same stem .pptx next to the PDF
        candidate = args.approval.with_suffix(".pptx")
        if candidate.is_file():
            pptx_twin = candidate
    if pptx_twin is not None:
        notes = try_pptx_notes(pptx_twin)
        if notes:
            approval_md += notes

    approval_out = OUT_DIR / "VaaniQ_Topic_Approval.md"
    approval_out.write_text(approval_md, encoding="utf-8")

    low_proposal = [p for p in proposal_pages if p["low_char"]]
    low_approval = [s for s in approval_slides if s["low_char"]]

    report: dict[str, Any] = {
        "ingested_at_utc": datetime.now(timezone.utc).isoformat(),
        "sources": {
            "proposal": {
                "path": str(args.proposal),
                "output": str(proposal_out.relative_to(ROOT)).replace("\\", "/"),
                "page_count": len(proposal_pages),
                "total_chars": sum(p["char_count"] for p in proposal_pages),
                "pages": proposal_pages,
                "low_char_pages": [p["page_number"] for p in low_proposal],
            },
            "topic_approval": {
                "path": str(args.approval),
                "output": str(approval_out.relative_to(ROOT)).replace("\\", "/"),
                "slide_count": len(approval_slides),
                "total_chars": sum(s["char_count"] for s in approval_slides),
                "slides": approval_slides,
                "low_char_slides": [s["slide_number"] for s in low_approval],
                "note": (
                    "Source provided as PDF export of the topic-approval presentation; "
                    "not a native PPTX. Slide markers used for citation parity with Phase 0."
                ),
            },
        },
        "summary": {
            "proposal_pages": len(proposal_pages),
            "approval_slides": len(approval_slides),
            "proposal_low_char_count": len(low_proposal),
            "approval_low_char_count": len(low_approval),
            "figures_rendered": len(low_proposal) + len(low_approval),
        },
    }

    report_path = OUT_DIR / "ingest_report.json"
    report_path.write_text(json.dumps(report, indent=2) + "\n", encoding="utf-8")

    print(json.dumps(report["summary"], indent=2))
    print(f"Wrote {proposal_out}")
    print(f"Wrote {approval_out}")
    print(f"Wrote {report_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
